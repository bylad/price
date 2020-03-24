import os
import sys
from bs4 import BeautifulSoup
import re
import requests
import xlrd
import docx
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import tkinter as tk
from tkinter import filedialog
from tkinter import messagebox

months = ['января', 'февраля', 'марта', 'апреля', 'мая', 'июня', 'июля',
          'августа', 'сентября', 'октября', 'ноября', 'декабря']


def dat(txt):
    """ Поиск даты по регулярному выражению.
        Например, из строки txt будет выделена дата:
        5 ноября 2019 года
    """
    dt = re.compile("\d+( |.)([яфмаисонд](([а-я]{2}\.? )|([а-я]+[а|я] ))\
                    d{4})|(\d+( |.)\d+)")
    match = re.search(dt, txt)
    if match:
        return (txt[match.start():match.end()])


def datastr(strlink):
    """ Вычленяем из даты числа, например, для записи в виде 20191105
        dig0-число, dig1-год, dig2-месяц
    """
    dig = [int(s) for s in strlink.split() if s.isdigit()]
    for string in months:
        regex = re.compile(string)
        match = re.search(regex, strlink)
        if match:
            text_pos = match.span()
            global month
            month = strlink[match.start():match.end()]
            monthnum = '{:02}'.format(months.index(month)+1)
            dig.append(int(monthnum))
    strlink_date = f'{dig[1]}-{monthnum[0]}-{dig[0]}'
    return(dig)


def news(newstext):
    """ Находим требуемую новость
    """
    a = soup.find('a', text=re.compile(newstext))
    return a


def newsin(file_name, textin):
    """ В найденной новости скачиваем файл с ценами НАО при его отсутствии,
        иначе без закачки приступаем к обработке файла.
    """
    if os.path.exists(file_name):
        print(f'Файл {file_name} существует, скачивание пропущено.')
        # print('Файл существует. Программа завершена!')
        # sys.exit()
    else:
        tovnao = ''
        rlink = requests.get(alink)
        soup = BeautifulSoup(rlink.text, 'html.parser')
        # for link in soup.findAll('a',text=re.compile('Ненецкому
        # автономному округу')):
        for link in soup.findAll('a', text=re.compile(textin)):
            tovnao = 'https://arhangelskstat.gks.ru'+link.get('href')
        # Добавляем headers, т.к. получали ошибку
        # "TimeoutError: [WinError 10060]" при попытке скачать файл
        # headers = {'user-agent':'Mozilla/5.0 (Windows NT 10.0; Win64; x64)
        # AppleWebKit/537.36 (KHTML, like Gecko) Chrome/62.0.3202.94
        # Safari/537.36'}
        rfile = requests.get(tovnao, headers=headers)
        with open(file_name, 'wb') as f:
            f.write(rfile.content)


def doc2docx(baseDir):
    """ Преобразуем файл DOC в DOCX
    """
    word = win32com.client.Dispatch("Word.application")
    for dir_path, dirs, files in os.walk(baseDir):
        for file_name in files:
            file_path = os.path.join(dir_path, file_name)
            file_name, file_extension = os.path.splitext(file_path)
            if file_extension.lower() == '.doc':
                docx_file = '{0}{1}'.format(file_path, 'x')
                # Skip conversion where docx file already exists
                if not os.path.isfile(docx_file):
                    print('Converting: {0}'.format(file_path))
                    try:
                        wordDoc = word.Documents.Open(file_path,
                                                      False, False, False)
                        wordDoc.SaveAs2(docx_file, FileFormat=16)
                        wordDoc.Close()
                    except Exception:
                        print('Failed to Convert: {0}'.format(file_path))
    # word.Quit()


def tabl(n, m):
    for j in range(1, n):
        try:
            cell_del_blank = re.sub(' ', '', table.cell(j, 1).text)
            cell = float(re.sub(',', '.', cell_del_blank))
        except IndexError:
            print(f'j={j}, j-1: {table.cell(j-1,0).text}=\
                  {table.cell(j-1,1).text}')
            sys.exit()

        table.cell(j, 1).text = str(t[j+m])
        if cell > t[j-1]:
            table.cell(j, 2).text = rdo
            table.cell(j, 2).text_frame.paragraphs[0].font.size = Pt(16)
            table.cell(j, 2).text_frame.paragraphs[0].font.bold = True
            table.cell(j, 2).text_frame.paragraphs[0].font.color.rgb =
            RGBColor(79, 98, 40)  # Green
        elif cell < t[j-1]:
            table.cell(j, 2).text = rup
            table.cell(j, 2).text_frame.paragraphs[0].font.size = Pt(16)
            table.cell(j, 2).text_frame.paragraphs[0].font.bold = True
            table.cell(j, 2).text_frame.paragraphs[0].font.color.rgb =
            RGBColor(158, 0, 0)  # Red
        else:
            table.cell(j, 2).text = ''
            table.cell(j, 2).text_frame.paragraphs[0].font.size = Pt(16)
        p = table.cell(j, 1).text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        p.font.name = 'Calibri'
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(37, 64, 97)


def pptx_path():
    """ Указываем путь к презентации со средними ценами по НАО
    """
    parent = tk.Tk()
    parent.withdraw()
    location = filedialog.askopenfilename(
        title="Выберите презентацию с ценами",
        filetypes=(("Powerpoint files", "*.pptx *.ppt"), ("all files", "*.*")))
    global prs = Presentation(location)
    directory = os.path.split(location)[0]
    # Поменять рабочий каталог на папку с файлом Powerpoint
    try:
        os.chdir(directory)
    except Exception:
        print("По указанному пути файл не может быть сохранен. \
              Скопируйте презентацию в доступную вам папку.")
        sys.exit()
    return directory

katalog = pptx_path()

url = 'https://arhangelskstat.gks.ru/news/'
# Добавляем headers, из-за ошибки при 1ом вызове alink = news()
# AttributeError: 'NoneType' object has no attribute 'get'
headers = {'user-agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) \
           AppleWebKit/537.36 (KHTML, like Gecko) Chrome/62.0.3202.94 \
           Safari/537.36'}
response = requests.get(url, headers=headers)
soup = BeautifulSoup(response.text, 'html.parser')

# Средние цены на отдельные потребительские товары
alink = news('Средние цены на отдельные потребительские товары').get('href')
digs = datastr(news('Средние цены на отдельные потребительские товары').text)
new_date = dat(news('Средние цены на отдельные потребительские товары').text)
datasite = f'{digs[1]}{digs[2]}{digs[0]}'
# print(f'Дата с сайта: {datasite}')
datasite_int = int(f'{digs[1]}{digs[2]}{digs[0]}')
# print(f'Дата в цифре: {datasite_int}')

newsin(fnamexls, 'Ненецкому автономному округу')
# -----------------------------------------------

# Цены на бензин АИ-92
alink = news('О потребительских ценах на бензин').get('href')
digsai = datastr(news('О потребительских ценах на бензин').text)
datasiteai = f'{digsai[1]}{digsai[2]}{digsai[0]}'
# print(f'Дата с сайта: {datasiteai}')
datasiteai_int = int(f'{digsai[1]}{digsai[2]}{digsai[0]}')

if datasite_int == datasiteai_int:
    fnamedoc = path2019+f'arhangelskstat_AI92_{datasiteai}.doc'  # DOC path
    newsin(fnamedoc, 'Ненецком автономном округе')
else:
    print('Внимание! НЕ совпадают даты новостей с ценами на бензин и \
          др. товары')
# -----------------------------------------------

t = []

# XLS файл со средними ценами на отдельные потребительские товары
path2019 = katalog + '/2019/'
fnamexls = path2019+f'arhangelskstat_{datasite}.xls'  # XLSX path
wb = xlrd.open_workbook(fnamexls)
sheet = wb.sheet_by_index(0)
'''
0   Говядина
1	Куры охлажденные
2	Колбаса вареная
3	Рыба мороженая
4	Молоко питьевое цельное пастеризованное
5	Яйца куриные
6	Чай черный байховый
7	Мука пшеничная
8	Хлеб из ржаной
9	Рис шлифованный
10	Картофель
11	Лук репчатый
12	Огурцы свежие
13	Яблоки
14	Майка
15	Мыло хозяйственное
16	Порошок
17	Проезд
18	Бензин автомобильный марки АИ-92
19	Плата за жилье
20	Отопление, Гкал
21	Водоснабжение холодное, м3
22	Водоотведение, м3
23	Водоснабжение горячее, м3
24	Услуги по снабжению
'''
tovar = ['Говядина', 'Куры охлажденные', 'Колбаса вареная',
         'Рыба мороженая', 'Молоко питьевое цельное пастеризованное',
         'Яйца куриные', 'Чай черный байховый', 'Мука пшеничная',
         'Хлеб из ржаной', 'Рис шлифованный', 'Картофель',
         'Лук репчатый', 'Огурцы свежие', 'Яблоки',
         'Майка', 'Мыло хозяйственное', 'Порошок',
         'Проезд', 'Бензин автомобильный марки АИ-92', 'Плата за жилье',
         'Отопление, Гкал', 'Водоснабжение холодное, м3', 'Водоотведение, м3',
         'Водоснабжение горячее, м3', 'Услуги по снабжению']

for i in range(sheet.nrows):
    for string in tovar:
        regex = re.compile(string)
        match = re.search(regex, sheet.cell_value(i, 0))
        if match:
            t.insert(tovar.index(string), sheet.cell_value(i, 1))
            # print(f'idx={tovar.index(string)},  str=\'{string}\',
            # {sheet.cell_value(i, 0)}={sheet.cell_value(i, 1)}')
# -----------------------------------------------

# DOC файл с ценами на бензин АИ-92
if not os.path.exists(fnamedoc+'x'):
    doc2docx(path2019)

doc = docx.Document(fnamedoc+'x')
ben = ['Бензин автомобильный марки АИ-92']
ri = 0
b = []
for table in doc.tables:
    for row in table.rows:
        for cell in row.cells:
            if any(x in cell.text for x in ben):
                flb = float(re.sub(',', '.', doc.tables[0].cell(ri, 1).text))
                b.append(flb)
                benz = doc.tables[0].cell(ri, 1).text
        ri = ri+1
# -----------------------------------------------
# Добавляем в наш список цену на бензин
t.insert(18, b[0])

# PPTX
prsoutf = pptx_path() + f'/midPriceOut_{datasite}.pptx'
rup = '↑'
rdo = '↓'
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.name == 'table0':
            table = shape.table
            tabl(15, -1)
        if shape.name == 'table1':
            table = shape.table
            tabl(4, 13)
        if shape.name == 'table2':
            table = shape.table
            tabl(9, 16)

        # Замена даты в заголовке
        if shape.name == 'naomp':
            text_frame = shape.text_frame
            cur_text = shape.text
            search_str = dat(shape.text)
            new_text = cur_text.replace(search_str, new_date)
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = new_text
            font = run.font
            font.name = 'Calibri'
            font.size = Pt(26)
            font.bold = True
            font.color.rgb = RGBColor(55, 96, 146)
print(f'Выходной файл лежит в папке:\n{prsoutf}')
prs.save(prsoutf)
